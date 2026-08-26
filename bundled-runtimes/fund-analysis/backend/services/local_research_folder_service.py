"""Read-only, incremental indexing for local research memo folders."""

from __future__ import annotations

import hashlib
import re
from copy import deepcopy
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional

from services.research_memo_viewpoint_taxonomy import ResearchMemoViewpointTaxonomy


class FolderValidationError(ValueError):
    """Raised when a folder cannot be indexed safely."""


class LocalResearchFolderService:
    SUPPORTED_SUFFIXES = {".md", ".txt", ".pdf", ".docx", ".pptx"}
    STYLE_LABELS = (
        "成长",
        "价值",
        "均衡",
        "质量",
        "红利",
        "大盘",
        "大中盘",
        "中盘",
        "中小盘",
        "小盘",
        "低换手",
        "高换手",
        "低波",
        "行业轮动",
        "主题",
        "周期",
        "景气",
        "量化",
        "指数增强",
        "固收+",
        "信用",
        "利率",
        "集中持仓",
        "分散持仓",
    )
    CLASSIFICATIONS = (
        "主动权益",
        "被动指数",
        "指数增强",
        "偏股混合",
        "灵活配置",
        "纯债",
        "一级债基",
        "二级债基",
        "货币基金",
        "QDII",
        "FOF",
        "商品",
        "固收",
        "行业主题",
        "周期资源",
    )
    STYLE_ALIASES = {
        "价值": (
            r"价值型", r"价值风格", r"偏价值", r"低估值", r"估值性价比",
            r"安全边际", r"左侧投资", r"逆向投资",
        ),
        "成长": (
            r"(?:大中盘|中小盘|大盘|中盘|小盘)?成长(?:型|风格|投资)?",
            r"偏成长",
        ),
        "均衡": (r"均衡型", r"均衡风格", r"风格均衡", r"均衡配置"),
        "质量": (r"质量型", r"质量价值", r"质量成长", r"高质量成长"),
        "红利": (r"红利(?:型|价值|风格|策略|投资|\+)", r"高股息", r"股息投资"),
        "大中盘": (r"大中盘",),
        "中小盘": (r"中小盘",),
        "大盘": (r"大盘型", r"大盘风格", r"大盘龙头"),
        "中盘": (r"中盘型", r"中盘风格"),
        "小盘": (r"小盘型", r"小盘风格", r"微盘"),
        "低换手": (r"低换手", r"换手率较低", r"降低换手率"),
        "高换手": (r"高换手", r"换手率较高"),
        "低波": (r"(?<!降)低波动", r"低波型", r"低波策略", r"波动率较低"),
        "行业轮动": (r"行业轮动",),
        "主题": (r"主题(?:型|基金|投资|策略|风格)", r"赛道基金", r"赛道型"),
        "周期": (r"周期(?:型|风格|投资|策略|思维|叠加|资源)", r"万物皆周期", r"看周期出身"),
        "景气": (r"景气(?:型|风格|投资|度投资|驱动|向上)",),
        "量化": (r"量化投资", r"主动量化", r"量化策略"),
        "指数增强": (r"指数增强",),
        "固收+": (r"固收\s*\+",),
        "信用": (r"信用债策略", r"信用策略"),
        "利率": (r"利率债策略", r"利率策略"),
        "集中持仓": (r"集中持仓", r"持仓集中", r"组合集中", r"重仓风格"),
        "分散持仓": (r"分散持仓", r"持仓分散", r"组合分散", r"行业和个股的分散"),
    }
    CLASSIFICATION_ALIASES = {
        "主动权益": (r"主动权益", r"主动股票", r"股票多头"),
        "被动指数": (r"被动指数", r"指数基金", r"ETF"),
        "指数增强": (r"指数增强",),
        "偏股混合": (r"偏股混合",),
        "灵活配置": (r"灵活配置",),
        "纯债": (r"纯债",),
        "一级债基": (r"一级债基",),
        "二级债基": (r"二级债基",),
        "货币基金": (r"货币基金", r"货币市场基金"),
        "QDII": (r"QDII",),
        "FOF": (r"FOF",),
        "商品": (r"商品基金", r"商品策略"),
        "固收": (r"固定收益", r"固收产品", r"固收策略"),
        "行业主题": (r"行业主题", r"主题基金", r"赛道基金"),
        "周期资源": (r"周期资源", r"资源周期", r"上游周期"),
    }
    METHOD_STYLE_LABELS = {
        "价值", "成长", "周期", "景气", "行业轮动", "量化", "指数增强",
        "固收+", "信用", "利率", "集中持仓", "分散持仓",
    }

    def __init__(
        self,
        repo: Any,
        manager_resolver: Optional[Callable[[str], Optional[Dict[str, Any]]]] = None,
        manager_fund_resolver: Optional[Callable[[str, str, str], List[Dict[str, Any]]]] = None,
        metadata_extractor: Optional[Callable[[str, str], Dict[str, Any]]] = None,
        manager_matcher: Optional[Any] = None,
        profile_projector: Optional[Callable[[Dict[str, Any], List[str]], Dict[str, Any]]] = None,
        manager_profile_projector: Optional[Callable[[Dict[str, Any], List[str]], Dict[str, Any]]] = None,
        max_files: int = 5_000,
        max_file_bytes: int = 25 * 1024 * 1024,
    ):
        self.repo = repo
        self.manager_resolver = manager_resolver
        self.manager_fund_resolver = manager_fund_resolver
        self.metadata_extractor = metadata_extractor
        self.manager_matcher = manager_matcher
        self.profile_projector = profile_projector
        self.manager_profile_projector = manager_profile_projector
        self.max_files = max_files
        self.max_file_bytes = max_file_bytes

    def add_folder(self, raw_path: str) -> Dict[str, Any]:
        path = self._validate_folder(raw_path)
        now = self._now()
        existing = next((item for item in self.repo.list_folders() if item.get("path") == str(path)), None)
        if existing:
            return existing
        return self.repo.create_folder({
            "path": str(path),
            "name": path.name,
            "status": "ready",
            "last_scan_at": None,
            "last_scan_counts": None,
            "created_at": now,
            "updated_at": now,
        })

    def list_folders(self) -> List[Dict[str, Any]]:
        return self.repo.list_folders()

    def scan_folder(self, folder_id: str, retry_llm: bool = False) -> Dict[str, Any]:
        folder = self.repo.get_folder(folder_id)
        if not folder:
            raise FolderValidationError("未找到已连接的调研文件夹")
        root = self._validate_folder(folder.get("path", ""))
        candidates = self._supported_files(root)
        if len(candidates) > self.max_files:
            raise FolderValidationError(f"可处理文件超过上限（{self.max_files} 份）")

        counts = {"created": 0, "updated": 0, "unchanged": 0, "failed": 0, "supported": len(candidates)}
        results: List[Dict[str, Any]] = []
        for path in candidates:
            result = self._index_file(folder_id, root, path, retry_llm=retry_llm)
            counts[result["status"]] += 1
            results.append(result)

        scanned_at = self._now()
        self.repo.update_folder(folder_id, {
            "status": "ready" if not counts["failed"] else "completed_with_errors",
            "last_scan_at": scanned_at,
            "last_scan_counts": counts,
            "updated_at": scanned_at,
        })
        projections = self._project_scan_results(results)
        return {
            "folder_id": folder_id,
            "folder_path": str(root),
            "scanned_at": scanned_at,
            "counts": counts,
            "results": results,
            "profile_projection": self._projection_summary(projections["fund_profile_projection"]),
            "fund_profile_projection": self._projection_summary(projections["fund_profile_projection"]),
            "manager_profile_projection": self._projection_summary(projections["manager_profile_projection"]),
        }

    def list_pending_reviews(self, folder_id: Optional[str] = None) -> List[Dict[str, Any]]:
        return self.repo.list_pending_reviews(folder_id)

    def confirm_manager_proposals(
        self,
        folder_id: Optional[str] = None,
        min_confidence: float = 0.88,
    ) -> Dict[str, Any]:
        candidates = [
            item for item in self.list_pending_reviews(folder_id)
            if item.get("kind") == "manager"
            and float(item.get("confidence") or 0) >= min_confidence
            and str(item.get("candidate_id") or "").strip()
            and str((item.get("identity_verification") or {}).get("status") or "") == "unique_exact_name"
            and str(item.get("report_date_source") or "") in {"filename", "content"}
            and self._has_explicit_manager_evidence(item)
        ]
        candidates_by_report: Dict[str, List[Dict[str, Any]]] = {}
        for item in candidates:
            report_id = str(item.get("report_id") or "").strip()
            if report_id:
                candidates_by_report.setdefault(report_id, []).append(item)

        confirmed = 0
        multi_manager = 0
        failed = 0
        linked_fund_count = 0
        for report_candidates in candidates_by_report.values():
            candidates_by_identity: Dict[str, List[Dict[str, Any]]] = {}
            for item in report_candidates:
                candidate_id = str(item.get("candidate_id") or "").strip()
                candidates_by_identity.setdefault(candidate_id, []).append(item)
            if len(candidates_by_identity) > 1:
                multi_manager += 1
            report_failed = False
            for identity_candidates in candidates_by_identity.values():
                item = max(
                    identity_candidates,
                    key=lambda candidate: float(candidate.get("confidence") or 0),
                )
                try:
                    result = self.review_proposal(
                        str(item.get("report_id") or ""),
                        str(item.get("id") or ""),
                        "confirmed",
                    )
                    linked_fund_count += int(result.get("linked_fund_count") or 0)
                except (TypeError, ValueError):
                    report_failed = True
            if report_failed:
                failed += 1
            else:
                confirmed += 1
        return {
            "status": "completed",
            "requested_reports": len(candidates_by_report),
            "confirmed": confirmed,
            "multi_manager": multi_manager,
            "ambiguous": 0,
            "failed": failed,
            "linked_fund_count": linked_fund_count,
            "min_confidence": min_confidence,
        }

    def confirm_label_proposals(
        self,
        folder_id: Optional[str] = None,
        min_confidence: float = 0.9,
    ) -> Dict[str, Any]:
        pending = [
            item for item in self.list_pending_reviews(folder_id)
            if item.get("kind") in {"classification", "style_label", "tag"}
            and float(item.get("confidence") or 0) >= min_confidence
            and item.get("extraction_source") != "llm"
        ]
        candidates = []
        skipped_unresolved = 0
        for item in pending:
            report = self.repo.get_report(str(item.get("report_id") or ""))
            if report and self._confirmable_label_proposal(item, report):
                candidates.append(item)
            else:
                skipped_unresolved += 1
        confirmed = 0
        failed = 0
        for item in candidates:
            try:
                self.review_proposal(
                    str(item.get("report_id") or ""),
                    str(item.get("id") or ""),
                    "confirmed",
                )
                confirmed += 1
            except (TypeError, ValueError):
                failed += 1
        return {
            "status": "completed",
            "requested": len(candidates),
            "confirmed": confirmed,
            "failed": failed,
            "skipped_unresolved": skipped_unresolved,
            "min_confidence": min_confidence,
        }

    def review_proposal(self, report_id: str, proposal_id: str, action: str) -> Dict[str, Any]:
        if action not in {"confirmed", "rejected"}:
            raise ValueError("复核结果只能是 confirmed 或 rejected")
        report = self.repo.get_report(report_id)
        if not report:
            raise ValueError("未找到调研纪要")
        old_fund_ids = list(report.get("fund_ids", []))
        old_manager_ids = self._report_manager_ids(report)

        proposals = report.get("review_proposals", [])
        target = next((item for item in proposals if item.get("id") == proposal_id), None)
        if not target:
            raise ValueError("未找到待复核项")
        target["review_status"] = action
        target["reviewed_at"] = self._now()

        fields = {
            "manager_id": report.get("manager_id") or "",
            "manager_name": report.get("manager_name") or "",
            "classifications": list(report.get("classifications", [])),
            "style_labels": list(report.get("style_labels", [])),
            "tags": list(report.get("tags", [])),
            "viewpoint_topics": list(report.get("viewpoint_topics", [])),
            "research_domains": list(report.get("research_domains", [])),
            "fund_ids": list(report.get("fund_ids", [])),
        }
        if target.get("kind") == "manager":
            manager_id = str(target.get("candidate_id") or "").strip()
            manager_name = str(target.get("value") or "").strip()
            if not manager_id:
                raise ValueError("基金经理缺少规范身份 ID")
            if action == "confirmed":
                self.repo.set_report_manager_link(
                    report_id,
                    manager_id,
                    manager_name,
                    source=str(target.get("extraction_source") or "research_memo_review"),
                    confirmed_at=target["reviewed_at"],
                )
            else:
                self.repo.remove_report_manager_link(report_id, manager_id)
            current_links = self.repo.list_report_manager_links(report_id)
            fields["manager_id"] = current_links[0]["manager_id"] if len(current_links) == 1 else ""
            fields["manager_name"] = current_links[0]["manager_name"] if len(current_links) == 1 else ""

        if target.get("kind") != "manager":
            self._apply_proposal(fields, target, confirmed=action == "confirmed")
        linked_fund_proposals = []
        if target.get("kind") == "manager":
            manager_name = str(target.get("value") or "").strip()
            for proposal in proposals:
                source_ref = proposal.get("source_ref") or {}
                if (
                    proposal.get("kind") != "fund"
                    or proposal.get("extraction_source") != "tushare.fund_manager"
                    or str(source_ref.get("manager_name") or "").strip() != manager_name
                ):
                    continue
                proposal["review_status"] = action
                proposal["reviewed_at"] = target["reviewed_at"]
                self._apply_proposal(fields, proposal, confirmed=action == "confirmed")
                linked_fund_proposals.append(proposal)
        fields.update({
            "review_proposals": proposals,
            "review_status": "pending" if any(
                item.get("review_status") == "pending" for item in proposals
            ) else "reviewed",
            "updated_at": self._now(),
        })
        updated = self.repo.update_report(report_id, fields)
        if updated:
            updated = self.repo.get_report(report_id) or updated
        affected_fund_ids = list(dict.fromkeys([
            *old_fund_ids,
            *(updated.get("fund_ids", []) if updated else []),
            *([target.get("value")] if target.get("kind") == "fund" else []),
        ]))
        affected_manager_ids = list(dict.fromkeys([
            *old_manager_ids,
            *self._report_manager_ids(updated or {}),
            *([str(target.get("candidate_id") or "").strip()] if target.get("kind") == "manager" else []),
        ]))
        affected_manager_ids = [item for item in affected_manager_ids if item]
        fund_projection = self.profile_projector(updated, affected_fund_ids) if self.profile_projector else None
        manager_projection = (
            self.manager_profile_projector(updated, affected_manager_ids)
            if self.manager_profile_projector else None
        )
        return {
            "status": action,
            "report": updated,
            "proposal": target,
            "linked_fund_proposals": linked_fund_proposals,
            "linked_fund_count": len(linked_fund_proposals),
            "profile_projection": fund_projection,
            "fund_profile_projection": fund_projection,
            "manager_profile_projection": manager_projection,
        }

    def _project_scan_results(self, results: List[Dict[str, Any]]) -> Dict[str, Any]:
        reports = []
        for report_id in dict.fromkeys(
            result.get("report_id")
            for result in results
            if result.get("report_id") and result.get("status") in {"created", "updated"}
        ):
            report = self.repo.get_report(report_id)
            if report:
                reports.append(report)
        fund_ids = list(dict.fromkeys(
            fund_id
            for report in reports
            for fund_id in report.get("fund_ids", [])
            if fund_id
        ))
        manager_ids = list(dict.fromkeys(
            manager_id
            for report in reports
            for manager_id in self._report_manager_ids(report)
        ))
        empty_fund_projection = {
                "projected_count": 0,
                "deleted_count": 0,
                "skipped_count": 0,
                "funds": [],
        }
        empty_manager_projection = {
            "projected_count": 0,
            "deleted_count": 0,
            "skipped_count": 0,
            "managers": [],
        }
        return {
            "fund_profile_projection": (
                self.profile_projector(reports[0], fund_ids)
                if self.profile_projector and reports and fund_ids else empty_fund_projection
            ),
            "manager_profile_projection": (
                self.manager_profile_projector(reports[0], manager_ids)
                if self.manager_profile_projector and reports and manager_ids else empty_manager_projection
            ),
        }

    def _index_file(
        self,
        folder_id: str,
        root: Path,
        path: Path,
        retry_llm: bool = False,
    ) -> Dict[str, Any]:
        relative_path = path.relative_to(root).as_posix()
        try:
            resolved_path = path.resolve(strict=True)
            if not resolved_path.is_relative_to(root):
                raise ValueError("文件指向已连接文件夹之外")
            stat = resolved_path.stat()
            existing = self.repo.get_document(folder_id, relative_path)
            if existing and existing.get("size") == stat.st_size and existing.get("mtime_ns") == stat.st_mtime_ns:
                enriched = self._enrich_existing_report(
                    existing.get("report_id"), root, resolved_path, retry_llm=retry_llm,
                )
                return {
                    "relative_path": relative_path,
                    "status": "updated" if enriched else "unchanged",
                    "report_id": existing.get("report_id"),
                }
            if stat.st_size > self.max_file_bytes:
                raise ValueError(f"文件超过 {self.max_file_bytes // (1024 * 1024) or 1} MB 上限")

            raw = resolved_path.read_bytes()
            content_hash = hashlib.sha256(raw).hexdigest()
            if existing and existing.get("content_hash") == content_hash:
                document = self._document_record(folder_id, root, resolved_path, stat, content_hash, existing.get("report_id"), "indexed")
                self.repo.upsert_document(document)
                enriched = self._enrich_existing_report(
                    existing.get("report_id"), root, resolved_path, retry_llm=retry_llm,
                )
                return {
                    "relative_path": relative_path,
                    "status": "updated" if enriched else "unchanged",
                    "report_id": existing.get("report_id"),
                }

            duplicate = self.repo.find_document_by_hash(
                content_hash,
                exclude_folder_id=folder_id,
                exclude_relative_path=relative_path,
            )
            if duplicate:
                document = self._document_record(folder_id, root, resolved_path, stat, content_hash, duplicate.get("report_id"), "duplicate")
                self.repo.upsert_document(document)
                return {"relative_path": relative_path, "status": "unchanged", "report_id": duplicate.get("report_id"), "duplicate": True}

            content = self._extract_text(resolved_path, raw).replace("\x00", "").strip()
            if not content:
                raise ValueError("未提取到可检索文字")
            extraction = self._extract_metadata(content, resolved_path.name)
            proposals = self._merge_proposals(
                self._extract_proposals(content, root, resolved_path),
                extraction.get("proposals", []),
                root,
                resolved_path,
            )
            now = self._now()
            report_date_evidence = self._report_date_evidence(resolved_path, stat.st_mtime, content)
            report_date = report_date_evidence["date"]
            proposals = self._add_manager_fund_proposals(
                proposals,
                report_date,
                resolved_path.stem,
                root,
                resolved_path,
            )
            proposals = self._scope_profile_proposals(proposals, content)
            viewpoint_topics = ResearchMemoViewpointTaxonomy.extract(content, resolved_path.stem)
            report_payload = {
                "manager_id": "",
                "manager_name": "",
                "title": resolved_path.stem,
                "report_date": report_date or None,
                "report_date_source": report_date_evidence["source"],
                "report_date_precision": report_date_evidence["precision"],
                "source": "本地调研纪要文件夹",
                "content": content,
                "summary": self._summary(content),
                "tags": [],
                "viewpoint_topics": viewpoint_topics,
                "research_domains": ResearchMemoViewpointTaxonomy.domains(
                    viewpoint_topics, content, resolved_path.stem
                ),
                "classifications": [],
                "style_labels": [],
                "fund_ids": [],
                "key_points": self._key_points(content),
                "review_proposals": proposals,
                "review_status": "pending" if proposals else "needs_metadata",
                "local_folder_id": folder_id,
                "local_relative_path": relative_path,
                "local_source_path": str(resolved_path),
                "source_hash": content_hash,
                "extraction_status": "complete",
                "extraction_provider": extraction.get("provider") or "deterministic_rules",
                "extraction_model": extraction.get("model"),
                "llm_extraction_status": extraction.get("status") or "unavailable",
                "llm_extraction_error": extraction.get("error"),
                "created_at": now,
                "updated_at": now,
            }
            if existing and existing.get("report_id"):
                report = self.repo.update_report(existing["report_id"], self._merge_review_state(report_payload, self.repo.get_report(existing["report_id"])))
                status = "updated"
            else:
                report = self.repo.create_report(report_payload)
                status = "created"
            document = self._document_record(folder_id, root, resolved_path, stat, content_hash, report["id"], "indexed")
            self.repo.upsert_document(document)
            return {"relative_path": relative_path, "status": status, "report_id": report["id"]}
        except Exception as exc:
            try:
                stat = path.stat()
                failed = {
                    "folder_id": folder_id,
                    "relative_path": relative_path,
                    "source_path": str(path),
                    "size": stat.st_size,
                    "mtime_ns": stat.st_mtime_ns,
                    "content_hash": None,
                    "report_id": None,
                    "index_status": "failed",
                    "error": str(exc),
                    "updated_at": self._now(),
                }
                self.repo.upsert_document(failed)
            except OSError:
                pass
            return {"relative_path": relative_path, "status": "failed", "error": str(exc)}

    def _extract_text(self, path: Path, raw: bytes) -> str:
        suffix = path.suffix.lower()
        if suffix in {".txt", ".md"}:
            for encoding in ("utf-8-sig", "gb18030"):
                try:
                    return raw.decode(encoding)
                except UnicodeDecodeError:
                    continue
            return raw.decode("utf-8", errors="replace")
        if suffix == ".pdf":
            from io import BytesIO
            from pypdf import PdfReader

            return "\n".join(page.extract_text() or "" for page in PdfReader(BytesIO(raw)).pages)
        if suffix == ".docx":
            from io import BytesIO
            from docx import Document

            document = Document(BytesIO(raw))
            return "\n".join(paragraph.text for paragraph in document.paragraphs)
        if suffix == ".pptx":
            from io import BytesIO
            from pptx import Presentation

            presentation = Presentation(BytesIO(raw))
            return "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            )
        raise ValueError(f"不支持的文件格式：{suffix}")

    def _extract_proposals(self, content: str, root: Path, path: Path) -> List[Dict[str, Any]]:
        proposals: List[Dict[str, Any]] = []
        if self.manager_matcher:
            for match in self.manager_matcher.match(content, path.name, path.relative_to(root).as_posix()):
                proposal = self._proposal(
                    "manager",
                    match["value"],
                    path,
                    root,
                    match.get("excerpt") or f"文件名：{path.name}",
                    float(match.get("confidence") or 0),
                    candidate_id=match.get("candidate_id"),
                )
                proposal["extraction_source"] = match.get("extraction_source") or "manager_matcher"
                if match.get("company"):
                    proposal["source_ref"]["candidate_company"] = match["company"]
                proposals.append(proposal)
        else:
            explicit_manager = re.search(r"基金经理\s*[：:]\s*([\u4e00-\u9fff·]{2,6})(?:先生|女士)?", content)
            if explicit_manager:
                value = re.sub(r"(?:先生|女士)$", "", explicit_manager.group(1).strip())
                candidate = self.manager_resolver(value) if self.manager_resolver else None
                proposals.append(self._proposal(
                    "manager", value, path, root,
                    self._line_excerpt(content, explicit_manager.start()),
                    0.98 if candidate else 0.88,
                    candidate_id=(candidate or {}).get("manager_id"),
                ))

        for match in re.finditer(r"(?<!\d)(\d{6}\.(?:OF|SH|SZ|BJ|HK))(?![A-Z])", content, re.IGNORECASE):
            value = match.group(1).upper()
            proposals.append(self._proposal(
                "fund",
                value,
                path,
                root,
                self._line_excerpt(content, match.start()),
                0.92,
            ))

        explicit_fields = (
            ("classification", r"(?:基金分类|基金类型|产品类型|产品定位)\s*[：:]\s*([^\r\n\f]+)"),
            (
                "style_label",
                r"(?:投资风格|基金风格|风格标签|持股风格|组合风格|产品风格|风格定位|(?<![\u4e00-\u9fff])风格)\s*[：:]\s*([^\r\n\f]+)",
            ),
        )
        for kind, pattern in explicit_fields:
            for match in re.finditer(pattern, content, re.IGNORECASE):
                stated_value = match.group(1).strip()
                values = (
                    self._explicit_style_values(stated_value)
                    if kind == "style_label"
                    else self._explicit_classification_values(stated_value)
                )
                for value in values:
                    proposal = self._proposal(
                        kind,
                        value,
                        path,
                        root,
                        re.sub(r"\s+", " ", match.group(0)).strip(),
                        0.94,
                    )
                    proposal["extraction_source"] = "explicit_field"
                    heading = re.sub(r"\s*[：:].*$", "", match.group(0)).strip()
                    proposal["source_ref"]["rule"] = (
                        "explicit_product_positioning"
                        if kind == "classification"
                        else "explicit_product_profile"
                        if heading in {"基金风格", "产品风格"}
                        else "explicit_manager_profile"
                    )
                    proposals.append(proposal)
        proposals.extend(self._deterministic_profile_proposals(content, path, root))
        return proposals

    @classmethod
    def _explicit_style_values(cls, stated_value: str) -> List[str]:
        """Extract only the profile clause, not market-style commentary."""
        clause = re.split(r"[。；;\n]", str(stated_value or "").strip(), maxsplit=1)[0].strip()
        if not clause or len(clause) > 160:
            return []
        if re.search(r"(?:跑赢|见顶|市场风格切换|风格轮动|阶段性行情|历史表现)", clause):
            return []
        values = [
            value
            for value, patterns in cls.STYLE_ALIASES.items()
            if cls._has_positive_pattern(clause, patterns)
        ]
        if "大中盘" in values:
            values = [value for value in values if value not in {"大盘", "中盘"}]
        if "中小盘" in values:
            values = [value for value in values if value not in {"中盘", "小盘"}]
        return list(dict.fromkeys(values))

    @classmethod
    def _explicit_classification_values(cls, stated_value: str) -> List[str]:
        clause = re.split(r"[。；;\n]", str(stated_value or "").strip(), maxsplit=1)[0].strip()
        if not clause or len(clause) > 160:
            return []
        return [
            value
            for value, patterns in cls.CLASSIFICATION_ALIASES.items()
            if cls._has_positive_pattern(clause, patterns)
        ]

    @staticmethod
    def _has_positive_pattern(text: str, patterns: tuple[str, ...]) -> bool:
        for pattern in patterns:
            for match in re.finditer(pattern, str(text or ""), re.IGNORECASE):
                prefix = str(text or "")[max(0, match.start() - 18):match.start()]
                if re.search(
                    r"(?:并非|不是|而不是|不做|不以|不限于|不局限于|避免|回避|不看好|不会|很少|未|非)\s*[^，。；;：:]{0,10}$",
                    prefix,
                ):
                    continue
                return True
        return False

    @staticmethod
    def _looks_like_profile_question(line: str) -> bool:
        text = re.sub(r"\s+", " ", str(line or "")).strip()
        return bool(
            re.match(r"^(?:Q\d*|问题\d*|提问|问)\s*[：:]", text, re.IGNORECASE)
            or "?" in text
            or "？" in text
            or re.search(r"(?:吗|呢)\s*[。.]?$", text)
            or re.search(r"(?:请|能否|可否).{0,24}(?:讲|介绍|分享)", text)
        )

    def _deterministic_profile_proposals(
        self,
        content: str,
        path: Path,
        root: Path,
    ) -> List[Dict[str, Any]]:
        """Extract only explicit self-described manager or product profile evidence."""
        proposals: List[Dict[str, Any]] = []
        explicit_context = re.compile(
            r"(?:我的|我自己|本人|个人).{0,18}(?:投资|持股|组合)?风格"
            r"|(?:产品|基金|投资|持股|组合)风格\s*(?:是|为|偏|定位|[:：])"
            r"|(?:投资框架|投资体系|投资方法|选股方法|组合构建|方法是)"
            r"|\d{6}\.(?:OF|SH|SZ|BJ|HK).{0,60}(?:风格|型)",
            re.IGNORECASE,
        )
        market_commentary = re.compile(
            r"(?:市场风格|风格切换|风格轮动|阶段性行情|历史表现|跑赢|见顶后|行情偏好)",
            re.IGNORECASE,
        )

        strongest: Dict[str, Dict[str, Any]] = {}
        classification_strongest: Dict[str, Dict[str, Any]] = {}
        for raw_line in str(content or "").splitlines():
            line = re.sub(r"\s+", " ", raw_line).strip()
            if not line or len(line) > 600 or self._looks_like_profile_question(line):
                continue
            if not explicit_context.search(line):
                continue
            if market_commentary.search(line) and not re.search(r"(?:我的|我自己|本人|个人|产品风格|组合风格|持股风格)", line):
                continue

            manager_style_match = re.search(
                r"(?:我的|我自己|本人|个人).{0,18}(?:投资|持股|组合)?风格"
                r"|(?:投资|持股|组合)风格\s*(?:是|为|偏|定位|[:：])",
                line,
                re.IGNORECASE,
            )
            product_style_match = re.search(
                r"(?:产品|基金)风格\s*(?:是|为|偏|定位|[:：])"
                r"|\d{6}\.(?:OF|SH|SZ|BJ|HK).{0,60}(?:风格|型)",
                line,
                re.IGNORECASE,
            )
            product_positioning_match = re.search(
                r"(?:内部)?定位为.{0,80}(?:基金|产品|策略)|(?:基金|产品).{0,24}(?:定位为|内部定位)",
                line,
                re.IGNORECASE,
            )
            is_manager_style = bool(manager_style_match)
            is_product_style = bool(product_style_match)
            is_product_positioning = bool(product_positioning_match)
            is_direct_style = is_manager_style or is_product_style
            label_text = line
            if is_direct_style:
                style_clause_match = re.search(
                    r"(?:我的|我自己|本人|个人).{0,18}(?:投资|持股|组合)?风格\s*(?:就是|是|为|偏|定位|[:：])?\s*(.+)"
                    r"|(?:产品|基金|投资|持股|组合)风格\s*(?:是|为|偏|定位|[:：])\s*(.+)"
                    r"|\d{6}\.(?:OF|SH|SZ|BJ|HK).{0,30}[：:]?\s*(.+?(?:风格|型).*)",
                    line,
                    re.IGNORECASE,
                )
                style_clause = next(
                    (group for group in (style_clause_match.groups() if style_clause_match else ()) if group),
                    line,
                )
                label_text = style_clause
            elif is_product_positioning and product_positioning_match:
                label_text = line[product_positioning_match.start():]
            labels = [
                label
                for label, patterns in self.STYLE_ALIASES.items()
                if self._has_positive_pattern(label_text, patterns)
            ]
            if re.search(r"(?:不集中|避免集中|降低集中度|控制集中度)", label_text):
                labels = [label for label in labels if label != "集中持仓"]
            if "大中盘" in labels:
                labels = [label for label in labels if label not in {"大盘", "中盘"}]
            if "中小盘" in labels:
                labels = [label for label in labels if label not in {"中盘", "小盘"}]
            if not is_direct_style and not is_product_positioning:
                labels = [label for label in labels if label in self.METHOD_STYLE_LABELS]
            confidence = 0.94 if is_direct_style else 0.91
            for label in labels:
                proposal = self._proposal("style_label", label, path, root, line, confidence)
                proposal["extraction_source"] = "deterministic_profile_rule"
                proposal["source_ref"]["rule"] = (
                    "explicit_manager_profile"
                    if is_manager_style
                    else "explicit_product_profile"
                    if is_product_style
                    else "explicit_product_positioning"
                    if is_product_positioning
                    else "investment_method_inference"
                )
                current = strongest.get(label)
                if not current or float(proposal["confidence"]) > float(current["confidence"]):
                    strongest[label] = proposal
            if is_product_positioning:
                for label, patterns in self.CLASSIFICATION_ALIASES.items():
                    if not self._has_positive_pattern(line, patterns):
                        continue
                    proposal = self._proposal("classification", label, path, root, line, 0.93)
                    proposal["extraction_source"] = "deterministic_profile_rule"
                    proposal["source_ref"]["rule"] = "explicit_product_positioning"
                    current = classification_strongest.get(label)
                    if not current or float(proposal["confidence"]) > float(current["confidence"]):
                        classification_strongest[label] = proposal
        proposals.extend(strongest.values())
        proposals.extend(classification_strongest.values())
        return proposals

    def _extract_metadata(self, content: str, filename: str) -> Dict[str, Any]:
        if not self.metadata_extractor:
            return {"status": "unavailable", "provider": None, "model": None, "proposals": []}
        try:
            result = self.metadata_extractor(content, filename)
            return result if isinstance(result, dict) else {"status": "failed", "proposals": [], "error": "模型提取结果格式无效"}
        except Exception as exc:
            return {"status": "failed", "provider": None, "model": None, "proposals": [], "error": str(exc)}

    def _enrich_existing_report(
        self,
        report_id: Optional[str],
        root: Path,
        path: Path,
        retry_llm: bool = False,
    ) -> bool:
        if not report_id:
            return False
        report = self.repo.get_report(report_id)
        if not report:
            return False
        source_path = Path(str(report.get("local_source_path") or path))
        if not source_path.exists():
            source_path = path
        proposal_path = source_path if source_path.is_relative_to(root) else path
        content = str(report.get("content") or "")
        if not content:
            content = self._extract_text(source_path, source_path.read_bytes()).strip()
        report_date_evidence = self._report_date_evidence(
            proposal_path,
            proposal_path.stat().st_mtime,
            content,
        )
        report_date = report_date_evidence["date"]
        previous_date = str(report.get("report_date") or "")
        previous_date_source = str(report.get("report_date_source") or "")
        if report_date_evidence["source"] == "unknown" and previous_date_source in {"filename", "content"}:
            report_date = previous_date
            report_date_evidence = {
                "date": previous_date,
                "source": previous_date_source,
                "precision": str(report.get("report_date_precision") or "day"),
            }
        extraction = None
        if retry_llm and report.get("llm_extraction_status") != "complete":
            extraction = self._extract_metadata(content, proposal_path.name)
        proposals = self._extract_proposals(content, root, proposal_path)
        proposals = self._merge_proposals(
            proposals,
            (extraction or {}).get("proposals") or [],
            root,
            proposal_path,
        )
        existing_llm_proposals = [
            proposal for proposal in (report.get("review_proposals") or [])
            if proposal.get("extraction_source") == "llm"
        ]
        proposals_by_key = {
            (proposal.get("kind"), proposal.get("value")): proposal
            for proposal in proposals
        }
        for proposal in existing_llm_proposals:
            key = (proposal.get("kind"), proposal.get("value"))
            current = proposals_by_key.get(key)
            if not current or float(proposal.get("confidence") or 0) > float(current.get("confidence") or 0):
                proposals_by_key[key] = proposal
        proposals = list(proposals_by_key.values())
        enriched = self._add_manager_fund_proposals(
            proposals,
            report_date,
            str(report.get("title") or proposal_path.stem),
            root,
            proposal_path,
        )
        enriched = self._scope_profile_proposals(enriched, content)
        merged_report = self._merge_review_state(
            {**report, "review_proposals": enriched},
            report,
        )
        enriched = merged_report["review_proposals"]
        viewpoint_topics = ResearchMemoViewpointTaxonomy.extract(
            content,
            str(report.get("title") or proposal_path.stem),
        )
        research_domains = ResearchMemoViewpointTaxonomy.domains(
            viewpoint_topics,
            content,
            str(report.get("title") or proposal_path.stem),
        )
        if (
            enriched == (report.get("review_proposals") or [])
            and report_date == str(report.get("report_date") or "")
            and report_date_evidence["source"] == str(report.get("report_date_source") or "")
            and report_date_evidence["precision"] == str(report.get("report_date_precision") or "")
            and viewpoint_topics == (report.get("viewpoint_topics") or [])
            and research_domains == (report.get("research_domains") or [])
            and extraction is None
        ):
            return False
        fields = {
            "report_date": report_date or None,
            "report_date_source": report_date_evidence["source"],
            "report_date_precision": report_date_evidence["precision"],
            "viewpoint_topics": viewpoint_topics,
            "research_domains": research_domains,
            "review_proposals": enriched,
            "manager_id": merged_report.get("manager_id") or "",
            "manager_name": merged_report.get("manager_name") or "",
            "classifications": merged_report.get("classifications") or [],
            "style_labels": merged_report.get("style_labels") or [],
            "tags": merged_report.get("tags") or [],
            "fund_ids": merged_report.get("fund_ids") or [],
            "review_status": "pending" if any(
                proposal.get("review_status") == "pending" for proposal in enriched
            ) else "reviewed",
            "updated_at": self._now(),
        }
        if extraction:
            fields.update({
                "extraction_provider": extraction.get("provider") or report.get("extraction_provider") or "deterministic_rules",
                "extraction_model": extraction.get("model"),
                "llm_extraction_status": extraction.get("status") or "unavailable",
                "llm_extraction_error": extraction.get("error"),
            })
        self.repo.update_report(report_id, fields)
        return True

    @staticmethod
    def _projection_summary(projection: Optional[Dict[str, Any]]) -> Dict[str, Any]:
        projection = projection or {}
        return {
            key: projection.get(key, default)
            for key, default in (
                ("updated_by", None),
                ("projected_count", 0),
                ("deleted_count", 0),
                ("cleared_count", 0),
                ("skipped_count", 0),
            )
        }

    def _add_manager_fund_proposals(
        self,
        proposals: List[Dict[str, Any]],
        report_date: str,
        report_title: str,
        root: Path,
        path: Path,
    ) -> List[Dict[str, Any]]:
        if not self.manager_fund_resolver:
            return proposals
        merged = {(item.get("kind"), item.get("value")): item for item in proposals}
        manager_names = list(dict.fromkeys(
            str(item.get("value") or "").strip()
            for item in proposals
            if item.get("kind") == "manager" and str(item.get("value") or "").strip()
        ))
        for manager_name in manager_names:
            try:
                relations = self.manager_fund_resolver(manager_name, report_date, report_title) or []
            except Exception:
                relations = []
            for relation in relations:
                wind_code = str(relation.get("wind_code") or "").strip().upper()
                if not wind_code:
                    continue
                proposal = self._proposal(
                    "fund",
                    wind_code,
                    path,
                    root,
                    f"Tushare：{manager_name} 于 {report_date} 管理 {relation.get('fund_name') or wind_code}",
                    0.9,
                )
                proposal["extraction_source"] = relation.get("source") or "tushare.fund_manager"
                proposal["source_ref"].update({
                    "manager_name": manager_name,
                    "report_date": report_date,
                    "fund_name": relation.get("fund_name"),
                    "management_company": relation.get("management_company"),
                })
                key = ("fund", wind_code)
                current = merged.get(key)
                if current and float(current.get("confidence") or 0) >= proposal["confidence"]:
                    current["extraction_source"] = proposal["extraction_source"]
                    current.setdefault("source_ref", {}).update({
                        "manager_name": manager_name,
                        "report_date": report_date,
                        "fund_name": relation.get("fund_name"),
                        "management_company": relation.get("management_company"),
                        "verification_excerpt": proposal["source_ref"]["excerpt"],
                    })
                else:
                    merged[key] = proposal
        return list(merged.values())

    def _scope_profile_proposals(
        self,
        proposals: List[Dict[str, Any]],
        content: str,
    ) -> List[Dict[str, Any]]:
        """Separate manager-level labels from labels tied to a named fund product."""
        fund_proposals = [item for item in proposals if item.get("kind") == "fund"]
        lines = [re.sub(r"\s+", " ", line).strip() for line in str(content or "").splitlines()]
        all_fund_references = list(dict.fromkeys(
            reference
            for fund in fund_proposals
            for reference in self._fund_references(fund)
        ))
        for proposal in proposals:
            if proposal.get("kind") not in {"classification", "style_label", "tag"}:
                continue
            value = str(proposal.get("value") or "").strip()
            target_fund_ids = [
                str(fund.get("value") or "").strip().upper()
                for fund in fund_proposals
                if self._label_targets_fund(lines, value, fund, all_fund_references)
            ]
            if (
                not target_fund_ids
                and len(fund_proposals) == 1
                and self._label_targets_unique_named_product(proposal, fund_proposals[0])
            ):
                target_fund_ids = [str(fund_proposals[0].get("value") or "").strip().upper()]
            proposal["target_fund_ids"] = list(dict.fromkeys(filter(None, target_fund_ids)))
            proposal["scope"] = "fund" if proposal["target_fund_ids"] else "manager"
        return proposals

    @classmethod
    def _label_targets_unique_named_product(
        cls,
        proposal: Dict[str, Any],
        fund_proposal: Dict[str, Any],
    ) -> bool:
        excerpt = re.sub(r"\s+", "", str((proposal.get("source_ref") or {}).get("excerpt") or ""))
        fund_name = re.sub(r"\s+", "", str((fund_proposal.get("source_ref") or {}).get("fund_name") or ""))
        if not excerpt or not fund_name or not re.search(r"(?:产品风格|产品定位|内部定位|定位为)", excerpt):
            return False
        base = re.sub(r"[-_/]?[A-H](?:类)?$", "", fund_name, flags=re.IGNORECASE)
        base = re.sub(r"(?:发起式)?(?:股票型?|混合型?|指数型?|债券型?|货币市场基金|基金)$", "", base)
        ignored = {"基金", "产品", "混合", "指数", "股票", "债券", "策略", "投资"}
        for size in range(min(len(base), 12), 3, -1):
            for start in range(0, len(base) - size + 1):
                fragment = base[start:start + size]
                if fragment in ignored or any(token in fragment for token in ("有限", "公司", "管理")):
                    continue
                if fragment in excerpt:
                    return True
        return False

    @classmethod
    def _label_targets_fund(
        cls,
        lines: List[str],
        label: str,
        fund_proposal: Dict[str, Any],
        all_fund_references: List[str],
    ) -> bool:
        references = cls._fund_references(fund_proposal)
        if not label or not references:
            return False

        style_cue = re.compile(r"(?:风格|策略|定位|偏好|指数增强|主动量化|红利策略|价值型|成长型|质量型)", re.IGNORECASE)
        for index, line in enumerate(lines):
            without_fund_names = line
            for reference in all_fund_references:
                without_fund_names = re.sub(re.escape(reference), "", without_fund_names, flags=re.IGNORECASE)
            if label.lower() not in without_fund_names.lower() or not style_cue.search(without_fund_names):
                continue
            if any(reference.lower() in line.lower() for reference in references):
                return True
            previous = lines[index - 1] if index > 0 else ""
            if re.search(r"(?:投资|产品|基金)?风格\s*[（(]", previous, re.IGNORECASE) and any(
                reference.lower() in previous.lower() for reference in references
            ):
                return True
        return False

    @classmethod
    def _fund_references(cls, fund_proposal: Dict[str, Any]) -> List[str]:
        code = str(fund_proposal.get("value") or "").strip().upper()
        source_ref = fund_proposal.get("source_ref") or {}
        return list(dict.fromkeys(filter(None, [
            code,
            code.split(".", 1)[0],
            *cls._fund_name_aliases(source_ref.get("fund_name")),
        ])))

    @staticmethod
    def _fund_name_aliases(fund_name: Any) -> List[str]:
        name = re.sub(r"\s+", "", str(fund_name or "")).strip()
        if not name:
            return []
        base = re.sub(r"[-_/]?[A-H](?:类)?$", "", name, flags=re.IGNORECASE)
        stripped = re.sub(
            r"(?:发起式)?(?:股票型?|混合型?|指数型?|债券型?|货币市场基金|基金)$",
            "",
            base,
        )
        return list(dict.fromkeys(alias for alias in (name, base, stripped) if len(alias) >= 5))

    def _merge_proposals(
        self,
        rule_proposals: List[Dict[str, Any]],
        model_proposals: List[Dict[str, Any]],
        root: Path,
        path: Path,
    ) -> List[Dict[str, Any]]:
        merged: Dict[tuple, Dict[str, Any]] = {}
        for item in rule_proposals:
            key = (item.get("kind"), item.get("value"))
            current = merged.get(key)
            if not current or float(item.get("confidence") or 0) > float(current.get("confidence") or 0):
                merged[key] = item
        for candidate in model_proposals:
            kind = candidate.get("kind")
            value = str(candidate.get("value") or "").strip()
            if kind not in {"manager", "fund", "classification", "style_label", "tag"} or not value:
                continue
            candidate_id = None
            if kind == "manager":
                excerpt = str(candidate.get("excerpt") or "")
                if self.manager_matcher and not self.manager_matcher.has_exact_name_evidence(value, excerpt):
                    continue
                normalized = None
                if self.manager_matcher:
                    normalized = self.manager_matcher.resolve_candidate(
                        value,
                        " ".join([
                            path.name,
                            path.relative_to(root).as_posix(),
                            str(candidate.get("excerpt") or ""),
                        ]),
                    )
                elif self.manager_resolver:
                    normalized = self.manager_resolver(value)
                if normalized:
                    value = str(normalized.get("value") or normalized.get("manager_name") or value).strip()
                    candidate_id = normalized.get("candidate_id") or normalized.get("manager_id")
                else:
                    continue
            key = (kind, value)
            proposal = self._proposal(
                kind,
                value,
                path,
                root,
                str(candidate.get("excerpt") or ""),
                float(candidate.get("confidence", 0)),
                candidate_id=candidate_id,
            )
            proposal["extraction_source"] = "llm"
            current = merged.get(key)
            if not current or proposal["confidence"] > current.get("confidence", 0):
                merged[key] = proposal
        return list(merged.values())

    @staticmethod
    def _has_explicit_manager_evidence(item: Dict[str, Any]) -> bool:
        source = str(item.get("extraction_source") or "")
        excerpt = str((item.get("source_ref") or {}).get("excerpt") or "").strip()
        manager_name = re.escape(str(item.get("value") or "").strip())
        if manager_name and re.search(
            rf"(?:历任|前任|原任|上一任|曾任)\s*(?:基金经理)?\s*[：:]?\s*{manager_name}",
            excerpt,
        ):
            return False
        return bool(
            excerpt
            and source in {"explicit_field", "manager_catalog_title", "filename_pattern", "llm"}
        )

    @staticmethod
    def _confirmable_label_proposal(item: Dict[str, Any], report: Dict[str, Any]) -> bool:
        excerpt = re.sub(r"\s+", " ", str((item.get("source_ref") or {}).get("excerpt") or "")).strip()
        value = str(item.get("value") or "").strip()
        if not excerpt or not value:
            return False
        source = str(item.get("extraction_source") or "")
        rule = str((item.get("source_ref") or {}).get("rule") or "")
        if source in {"deterministic_profile_rule", "explicit_field"} and rule not in {
            "explicit_manager_profile",
            "explicit_product_profile",
            "explicit_product_positioning",
        }:
            return False
        scope = str(item.get("scope") or "manager")
        if scope == "fund":
            targets = {
                str(code or "").strip().upper()
                for code in item.get("target_fund_ids") or []
                if str(code or "").strip()
            }
            report_funds = {
                str(code or "").strip().upper()
                for code in report.get("fund_ids") or []
                if str(code or "").strip()
            }
            return bool(len(targets) == 1 and targets.issubset(report_funds))
        return bool(
            item.get("kind") == "style_label"
            and rule == "explicit_manager_profile"
            and LocalResearchFolderService._report_manager_ids(report)
        )

    def _proposal(
        self,
        kind: str,
        value: str,
        path: Path,
        root: Path,
        excerpt: str,
        confidence: float,
        candidate_id: Optional[str] = None,
    ) -> Dict[str, Any]:
        relative_path = path.relative_to(root).as_posix()
        identity = hashlib.sha256(f"{relative_path}\0{kind}\0{value}".encode("utf-8")).hexdigest()[:20]
        return {
            "id": identity,
            "kind": kind,
            "value": value,
            "candidate_id": candidate_id,
            "confidence": confidence,
            "review_status": "pending",
            "source_ref": {
                "relative_path": relative_path,
                "source_path": str(path),
                "excerpt": excerpt[:240],
            },
        }

    @staticmethod
    def _apply_proposal(fields: Dict[str, Any], proposal: Dict[str, Any], confirmed: bool) -> None:
        kind = proposal.get("kind")
        value = proposal.get("value")
        if kind == "manager":
            # 正式关联写入 research_report_managers；旧字段只兼容恰好一位经理的纪要。
            if confirmed:
                current_id = str(fields.get("manager_id") or "").strip()
                candidate_id = str(proposal.get("candidate_id") or "").strip()
                if not current_id or current_id == candidate_id:
                    fields["manager_name"] = value
                    fields["manager_id"] = candidate_id
                else:
                    fields["manager_name"] = ""
                    fields["manager_id"] = ""
            elif str(fields.get("manager_id") or "").strip() == str(proposal.get("candidate_id") or "").strip():
                fields["manager_name"] = ""
                fields["manager_id"] = ""
            return
        field_name = {
            "fund": "fund_ids",
            "classification": "classifications",
            "style_label": "style_labels",
            "tag": "tags",
        }.get(kind)
        if not field_name:
            return
        values = [item for item in fields.get(field_name, []) if item != value]
        if confirmed:
            values.append(value)
        fields[field_name] = list(dict.fromkeys(values))

    def _merge_review_state(self, fresh: Dict[str, Any], existing: Optional[Dict[str, Any]]) -> Dict[str, Any]:
        if not existing:
            return fresh
        fresh_by_key = {
            (item.get("kind"), item.get("value")): item
            for item in fresh.get("review_proposals", [])
        }
        for old in existing.get("review_proposals", []):
            key = (old.get("kind"), old.get("value"))
            if key not in fresh_by_key:
                retained = deepcopy(old)
                fresh["review_proposals"].append(retained)
                fresh_by_key[key] = retained
        old_by_key = {
            (item.get("kind"), item.get("value")): item
            for item in existing.get("review_proposals", [])
            if item.get("review_status") in {"confirmed", "rejected"}
        }
        fields = {
            "manager_id": "",
            "manager_name": "",
            "classifications": [],
            "style_labels": [],
            "tags": [],
            "viewpoint_topics": list(fresh.get("viewpoint_topics") or existing.get("viewpoint_topics") or []),
            "research_domains": list(fresh.get("research_domains") or existing.get("research_domains") or []),
            "fund_ids": [],
        }
        for proposal in fresh["review_proposals"]:
            old = old_by_key.get((proposal.get("kind"), proposal.get("value")))
            if old:
                proposal["review_status"] = old["review_status"]
                proposal["reviewed_at"] = old.get("reviewed_at")
                proposal["candidate_id"] = old.get("candidate_id") or proposal.get("candidate_id")
                if old.get("identity_verification"):
                    proposal["identity_verification"] = old["identity_verification"]

        manager_links = list(existing.get("manager_links") or [])
        links_by_id = {
            str(link.get("manager_id") or "").strip(): link
            for link in manager_links
            if str(link.get("manager_id") or "").strip()
        }
        for proposal in fresh["review_proposals"]:
            if proposal.get("kind") != "manager" or proposal.get("review_status") != "pending":
                continue
            manager_id = str(proposal.get("candidate_id") or "").strip()
            link = links_by_id.get(manager_id)
            if not link:
                continue
            proposal["review_status"] = "confirmed"
            proposal["reviewed_at"] = link.get("confirmed_at") or self._now()

        confirmed_managers = {
            str(proposal.get("value") or "").strip(): proposal.get("reviewed_at")
            for proposal in fresh["review_proposals"]
            if proposal.get("kind") == "manager" and proposal.get("review_status") == "confirmed"
        }
        for proposal in fresh["review_proposals"]:
            source_ref = proposal.get("source_ref") or {}
            manager_name = str(source_ref.get("manager_name") or "").strip()
            if (
                proposal.get("kind") == "fund"
                and proposal.get("review_status") == "pending"
                and proposal.get("extraction_source") == "tushare.fund_manager"
                and manager_name in confirmed_managers
            ):
                proposal["review_status"] = "confirmed"
                proposal["reviewed_at"] = confirmed_managers[manager_name] or self._now()

        for proposal in fresh["review_proposals"]:
            if proposal.get("review_status") == "confirmed" and proposal.get("kind") != "manager":
                self._apply_proposal(fields, proposal, confirmed=True)
        if len(manager_links) == 1:
            fields["manager_id"] = manager_links[0].get("manager_id") or ""
            fields["manager_name"] = manager_links[0].get("manager_name") or ""
        fresh.update(fields)
        fresh["created_at"] = existing.get("created_at") or fresh["created_at"]
        fresh["review_status"] = "pending" if any(
            item.get("review_status") == "pending" for item in fresh["review_proposals"]
        ) else "reviewed"
        return fresh

    @staticmethod
    def _report_manager_ids(report: Dict[str, Any]) -> List[str]:
        manager_ids = [
            str(item.get("manager_id") or "").strip()
            for item in report.get("manager_links") or []
            if str(item.get("manager_id") or "").strip()
        ]
        legacy_manager_id = str(report.get("manager_id") or "").strip()
        if legacy_manager_id:
            manager_ids.append(legacy_manager_id)
        return list(dict.fromkeys(manager_ids))

    def _supported_files(self, root: Path) -> List[Path]:
        files = []
        for path in root.rglob("*"):
            if path.is_symlink() or not path.is_file() or path.suffix.lower() not in self.SUPPORTED_SUFFIXES:
                continue
            files.append(path)
        return sorted(files, key=lambda item: item.relative_to(root).as_posix().casefold())

    def _validate_folder(self, raw_path: str) -> Path:
        if not str(raw_path or "").strip():
            raise FolderValidationError("请输入调研纪要文件夹路径")
        try:
            path = Path(raw_path).expanduser().resolve(strict=True)
        except (OSError, RuntimeError) as exc:
            raise FolderValidationError("文件夹不存在或无法读取") from exc
        if not path.is_dir():
            raise FolderValidationError("所选路径不是文件夹")
        if path in {Path(path.anchor), Path.home().resolve()}:
            raise FolderValidationError("不能连接系统根目录或整个用户目录")
        return path

    def _document_record(
        self,
        folder_id: str,
        root: Path,
        path: Path,
        stat: Any,
        content_hash: str,
        report_id: str,
        index_status: str,
    ) -> Dict[str, Any]:
        return {
            "folder_id": folder_id,
            "relative_path": path.relative_to(root).as_posix(),
            "source_path": str(path),
            "size": stat.st_size,
            "mtime_ns": stat.st_mtime_ns,
            "content_hash": content_hash,
            "report_id": report_id,
            "index_status": index_status,
            "error": None,
            "updated_at": self._now(),
        }

    @staticmethod
    def _summary(content: str) -> str:
        return re.sub(r"\s+", " ", content).strip()[:500]

    @staticmethod
    def _key_points(content: str) -> List[str]:
        points = []
        for line in content.splitlines():
            clean = line.strip().lstrip("-•* ").strip()
            if line.strip().startswith(("-", "•", "*")) and 8 <= len(clean) <= 180:
                points.append(clean)
        return points[:5]

    @staticmethod
    def _report_date(path: Path, modified_at: float, content: str = "") -> str:
        return LocalResearchFolderService._report_date_evidence(path, modified_at, content)["date"]

    @staticmethod
    def _report_date_evidence(path: Path, modified_at: float, content: str = "") -> Dict[str, str]:
        name = path.stem
        date_patterns = (
            r"(?<!\d)(20\d{2})[年._-](\d{1,2})[月._-](\d{1,2})日?(?!\d)",
            r"(?<!\d)(20\d{2})(\d{2})(\d{2})(?!\d)",
            r"(?<!\d)(\d{2})年(\d{1,2})月(\d{1,2})日(?!\d)",
            r"(?<!\d)(\d{2})(\d{2})(\d{2})(?!\d)",
        )
        for pattern in date_patterns:
            match = re.search(pattern, name)
            if not match:
                continue
            year = int(match.group(1))
            if year < 100:
                year += 2000
            try:
                return {
                    "date": datetime(year, int(match.group(2)), int(match.group(3)), tzinfo=timezone.utc).date().isoformat(),
                    "source": "filename",
                    "precision": "day",
                }
            except ValueError:
                continue
        compact_month = re.search(r"(?<!\d)(20\d{2})(0[1-9]|1[0-2])(?!\d)", name)
        if compact_month:
            return {
                "date": datetime(
                    int(compact_month.group(1)),
                    int(compact_month.group(2)),
                    1,
                    tzinfo=timezone.utc,
                ).date().isoformat(),
                "source": "filename",
                "precision": "month",
            }
        compact_short_month = re.search(r"(?<!\d)(\d{2})(0[1-9]|1[0-2])(?!\d)", name)
        if compact_short_month:
            year = 2000 + int(compact_short_month.group(1))
            path_years = {
                int(part)
                for part in path.parts
                if re.fullmatch(r"20\d{2}", part)
            }
            if year in path_years:
                return {
                    "date": datetime(
                        year,
                        int(compact_short_month.group(2)),
                        1,
                        tzinfo=timezone.utc,
                    ).date().isoformat(),
                    "source": "filename",
                    "precision": "month",
                }
        quarter = re.search(r"(?<!\d)(20)?(\d{2})\s*[Qq]([1-4])(?!\d)", name)
        if quarter:
            return {
                "date": datetime(
                    2000 + int(quarter.group(2)),
                    (int(quarter.group(3)) - 1) * 3 + 1,
                    1,
                    tzinfo=timezone.utc,
                ).date().isoformat(),
                "source": "filename",
                "precision": "quarter",
            }
        content_lines = [re.sub(r"\s+", " ", line).strip() for line in str(content or "").splitlines()[:40]]
        explicit_date = next((
            match
            for line in content_lines
            if not re.search(r"(?:数据|业绩|任职回报|截至|截止|来源)", line)
            if (match := re.search(
                r"(?:路演|会议|交流|调研|访谈|报告|材料)?(?:时间|日期)\s*[：:]?\s*"
                r"(20\d{2})[年./_-](\d{1,2})[月./_-](\d{1,2})日?",
                line,
            ))
        ), None)
        if explicit_date:
            try:
                return {
                    "date": datetime(
                        int(explicit_date.group(1)),
                        int(explicit_date.group(2)),
                        int(explicit_date.group(3)),
                        tzinfo=timezone.utc,
                    ).date().isoformat(),
                    "source": "content",
                    "precision": "day",
                }
            except ValueError:
                pass
        standalone_date = next((
            match
            for line in content_lines[:12]
            if (match := re.fullmatch(r"(20\d{2})年(\d{1,2})月(?:(\d{1,2})日)?", line))
        ), None)
        if standalone_date:
            try:
                return {
                    "date": datetime(
                        int(standalone_date.group(1)),
                        int(standalone_date.group(2)),
                        int(standalone_date.group(3) or 1),
                        tzinfo=timezone.utc,
                    ).date().isoformat(),
                    "source": "content",
                    "precision": "day" if standalone_date.group(3) else "month",
                }
            except ValueError:
                pass
        return {
            "date": "",
            "source": "unknown",
            "precision": "unknown",
        }

    @staticmethod
    def _line_excerpt(content: str, character_index: int) -> str:
        start = content.rfind("\n", 0, character_index) + 1
        end = content.find("\n", character_index)
        return content[start: end if end >= 0 else len(content)].strip()

    @classmethod
    def _excerpt_for_value(cls, content: str, value: str) -> str:
        index = content.lower().find(value.lower())
        return cls._line_excerpt(content, max(index, 0))

    @staticmethod
    def _now() -> str:
        return datetime.now(timezone.utc).isoformat()
